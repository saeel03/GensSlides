from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
import os, random, re, uuid
import g4f
from pdf2docx import Converter
from docx2pdf import convert as docx2pdf_convert  # ✅ For Word to PDF
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader


app = Flask(__name__)
CORS(app)

DESIGNS_DIR = "backend/Designs"
OUTPUT_DIR = "backend/generated_ppt"
UPLOAD_FOLDER = "backend/uploads"
CONVERTED_FOLDER = "backend/converted"

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

def clean_text_lines(text: str) -> list[str]:
    return [line.strip() for line in re.split(r'\n+', text) if line.strip()]

def generate_slide_content(topic: str) -> str:
    prompt = (
        f"Create a detailed PowerPoint presentation content outline for the topic '{topic}'.\n"
        "Each slide should have a clear and meaningful title, followed by 3–5 detailed bullet points.\n"
        "Format:\n## Slide 1: Slide Title\n- Bullet 1\n- Bullet 2\n"
        "Include intro, background, detailed explanation, examples, conclusion, and Q&A."
    )
    try:
        response = g4f.ChatCompletion.create(
            model=g4f.models.default,
            messages=[{"role": "user", "content": prompt}],
        )
        print("✅ AI response received.")
        return response
    except Exception as e:
        print(f"⚠️ AI generation failed: {e}")
        return "## Slide 1: Introduction\n- Overview\n## Slide 2: Key Concepts\n- Explanation"

def build_presentation(content: str, template_path: str, topic: str):
    try:
        prs = Presentation(template_path)

        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = topic
        title_slide.placeholders[1].text = f"Exploring the Future of {topic}"

        content_layout = prs.slide_layouts[1]
        slide_blocks = re.findall(r'##\s*Slide\s*\d+:\s*(.+?)\n((?:- .+\n?)*)', content, re.MULTILINE)
        if not slide_blocks:
            print("⚠️ No slide blocks matched.")
            return None, None

        for slide_title, bullets in slide_blocks:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_title.strip()
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for point in clean_text_lines(bullets):
                p = text_frame.add_paragraph()
                p.level = 0
                p.text = point

        filename = f"{uuid.uuid4()}.pptx"
        output_path = os.path.join(OUTPUT_DIR, filename)
        prs.save(output_path)
        return output_path, filename
    except Exception as e:
        print(f"❌ Error: {e}")
        return None, None

@app.route('/generate-ppt-topic', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()
        topic = data.get('topic', '').strip()
        if not topic:
            return jsonify({"error": "No topic provided"}), 400

        content = generate_slide_content(topic)

        templates = [os.path.join(DESIGNS_DIR, f) for f in os.listdir(DESIGNS_DIR) if f.endswith('.pptx')]
        if not templates:
            return jsonify({"error": "No templates found"}), 500

        selected_template = random.choice(templates)
        ppt_path, filename = build_presentation(content, selected_template, topic)
        if not ppt_path:
            return jsonify({"error": "PPT generation failed"}), 500

        return jsonify({"message": "Presentation generated", "download_url": f"http://localhost:5000/download/{filename}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/generate-ppt-from-content', methods=['POST'])
def generate_ppt_from_content():
    try:
        data = request.get_json()
        content_input = data.get('content', '').strip()
        if not content_input:
            return jsonify({"error": "No content provided"}), 400

        # Summarize and convert content to slide-friendly format with short titles
        summarization_prompt = (
            f"Summarize the following content into a PowerPoint presentation structure.\n\n"
            f"Ensure each slide has a **short, clear title** (1-5 words) and 3–5 detailed bullet points.\n\n"
            f"{content_input}\n\n"
            "Format:\n## Slide 1: Short Title\n- Bullet 1\n- Bullet 2\n..."
        )

        ai_response = g4f.ChatCompletion.create(
            model=g4f.models.default,
            messages=[{"role": "user", "content": summarization_prompt}],
        )
        

        print("✅ AI summary for content received.")

        # Generate a short title for the overall presentation
        title_prompt = (
            f"Generate a short and relevant title (max 7 words) for a presentation based on the following content:\n\n{content_input}"
        )
        title_response = g4f.ChatCompletion.create(
            model=g4f.models.default,
            messages=[{"role": "user", "content": title_prompt}],
        )
        cleaned_title = title_response.strip().replace('"', '')
        if not cleaned_title or "summarized content" in cleaned_title.lower():
            cleaned_title = "AI Generated Presentation"

        templates = [os.path.join(DESIGNS_DIR, f) for f in os.listdir(DESIGNS_DIR) if f.endswith('.pptx')]
        if not templates:
            return jsonify({"error": "No templates found"}), 500

        selected_template = random.choice(templates)
        ppt_path, filename = build_presentation(ai_response, selected_template, topic=cleaned_title)

        if not ppt_path:
            return jsonify({"error": "PPT generation failed"}), 500

        return jsonify({
            "message": "Presentation generated from content",
            "download_url": f"http://localhost:5000/download/{filename}"
        })

    except Exception as e:
        print(f"❌ Exception in content-to-PPT: {e}")
        return jsonify({"error": str(e)}), 500





@app.route('/download/<filename>')
def download_ppt(filename):
    path = os.path.join(OUTPUT_DIR, filename)
    return send_file(path, as_attachment=True) if os.path.exists(path) else ("File not found", 404)

@app.route('/pdf-to-word', methods=['POST'])
def convert_pdf_to_word():
    if 'pdf' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['pdf']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    filename = secure_filename(file.filename)
    pdf_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(pdf_path)

    word_filename = filename.rsplit('.', 1)[0] + '.docx'
    word_path = os.path.join(CONVERTED_FOLDER, word_filename)

    try:
        cv = Converter(pdf_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()
    except Exception as e:
        return jsonify({'error': str(e)}), 500

    return jsonify({'download_url': f"http://localhost:5000/converted/{word_filename}"})

@app.route('/word-to-pdf', methods=['POST'])
def convert_word_to_pdf():
    if 'word' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['word']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    filename = secure_filename(file.filename)
    word_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(word_path)

    pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
    pdf_path = os.path.join(CONVERTED_FOLDER, pdf_filename)

    try:
        docx2pdf_convert(word_path, pdf_path)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

    return jsonify({'download_url': f"http://localhost:5000/converted/{pdf_filename}"})

@app.route('/pdf-to-ppt', methods=['POST'])
def pdf_to_ppt():
    if 'pdf' not in request.files:
        return jsonify({'error': 'No PDF uploaded'}), 400

    file = request.files['pdf']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    filename = secure_filename(file.filename)
    pdf_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(pdf_path)

    try:
        # Extract text from PDF
        reader = PdfReader(pdf_path)
        content = ''
        for page in reader.pages:
            content += page.extract_text() or ""

        if not content.strip():
            return jsonify({'error': 'Unable to extract text from PDF'}), 400

        # Summarize into slides
        summarization_prompt = (
            f"Summarize the following content into a PowerPoint presentation structure.\n\n"
            f"Each slide should have a short, clear title (1–5 words) and 3–5 bullet points.\n\n"
            f"{content}\n\n"
            "Format:\n## Slide 1: Short Title\n- Bullet 1\n- Bullet 2\n..."
        )
        summary_response = g4f.ChatCompletion.create(
            model=g4f.models.default,
            messages=[{"role": "user", "content": summarization_prompt}]
        )

        # Get presentation title
        title_prompt = f"Generate a short, relevant title (max 7 words) for a presentation based on:\n\n{content}"
        title_response = g4f.ChatCompletion.create(
            model=g4f.models.default,
            messages=[{"role": "user", "content": title_prompt}]
        )
        title = title_response.strip().replace('"', '')
        if not title or "summary" in title.lower():
            title = "PDF Based Presentation"

        # Pick a random design template
        templates = [os.path.join(DESIGNS_DIR, f) for f in os.listdir(DESIGNS_DIR) if f.endswith('.pptx')]
        if not templates:
            return jsonify({"error": "No templates found"}), 500

        selected_template = random.choice(templates)
        ppt_path, filename = build_presentation(summary_response, selected_template, topic=title)

        if not ppt_path:
            return jsonify({'error': 'PPT generation failed'}), 500

        return jsonify({'message': 'PPT created from PDF', 'download_url': f'http://localhost:5000/download/{filename}'})

    except Exception as e:
        print(f"❌ PDF to PPT error: {e}")
        return jsonify({'error': str(e)}), 500
    
@app.route('/converted/<filename>')
def download_converted_file(filename):
    file_path = os.path.join(CONVERTED_FOLDER, filename)
    return send_file(file_path, as_attachment=True) if os.path.exists(file_path) else ("File not found", 404)

if __name__ == '__main__':
    app.run(debug=True)
